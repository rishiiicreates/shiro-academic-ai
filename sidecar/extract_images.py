import os
import sys
import json
import re
import io
import pymupdf
import pptx
from pptx.enum.shapes import MSO_SHAPE_TYPE

SOURCE_DIR = '/Users/rishii/THE-Helper'
OUTPUT_DIR = '/Users/rishii/the-helper-rag-app/data/images'
MANIFEST_PATH = '/Users/rishii/the-helper-rag-app/data/images_manifest.json'

os.makedirs(OUTPUT_DIR, exist_ok=True)

def sanitize_filename(name):
    # Keep alphanumeric, dot, underscore, hyphen
    clean = re.sub(r'[^\w\-_\.]', '_', name)
    return clean[:120]

def extract_pdf_images(file_path, file_name, manifest):
    try:
        doc = pymupdf.open(file_path)
    except Exception as e:
        print(f"Error opening PDF {file_name}: {e}")
        return 0

    total_extracted = 0
    safe_dir_name = sanitize_filename(file_name)
    file_img_dir = os.path.join(OUTPUT_DIR, safe_dir_name)

    for page_idx, page in enumerate(doc):
        page_num = str(page_idx + 1)
        image_list = page.get_images(full=True)
        if not image_list:
            continue

        for img_idx, img_info in enumerate(image_list):
            xref = img_info[0]
            try:
                base_img = doc.extract_image(xref)
                if not base_img:
                    continue
                
                width = base_img.get("width", 0)
                height = base_img.get("height", 0)
                ext = base_img.get("ext", "png")
                img_bytes = base_img.get("image", b"")

                # Filter out tiny bullets, icons, transparent pixels
                if width < 120 or height < 120 or len(img_bytes) < 3000:
                    continue

                os.makedirs(file_img_dir, exist_ok=True)
                img_filename = f"p{page_num}_{img_idx}.{ext}"
                out_path = os.path.join(file_img_dir, img_filename)
                
                with open(out_path, "wb") as f:
                    f.write(img_bytes)

                rel_url = f"/images/{safe_dir_name}/{img_filename}"
                
                if file_name not in manifest:
                    manifest[file_name] = {}
                if page_num not in manifest[file_name]:
                    manifest[file_name][page_num] = []
                
                if rel_url not in manifest[file_name][page_num]:
                    manifest[file_name][page_num].append(rel_url)
                
                total_extracted += 1
            except Exception as ex:
                continue

    return total_extracted

def extract_pptx_images(file_path, file_name, manifest):
    try:
        prs = pptx.Presentation(file_path)
    except Exception as e:
        print(f"Error opening PPTX {file_name}: {e}")
        return 0

    total_extracted = 0
    safe_dir_name = sanitize_filename(file_name)
    file_img_dir = os.path.join(OUTPUT_DIR, safe_dir_name)

    for slide_idx, slide in enumerate(prs.slides):
        page_num = str(slide_idx + 1)
        img_idx = 0

        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image = shape.image
                    img_bytes = image.blob
                    ext = image.ext or "png"

                    if len(img_bytes) < 3000:
                        continue

                    os.makedirs(file_img_dir, exist_ok=True)
                    img_filename = f"p{page_num}_{img_idx}.{ext}"
                    out_path = os.path.join(file_img_dir, img_filename)

                    with open(out_path, "wb") as f:
                        f.write(img_bytes)

                    rel_url = f"/images/{safe_dir_name}/{img_filename}"

                    if file_name not in manifest:
                        manifest[file_name] = {}
                    if page_num not in manifest[file_name]:
                        manifest[file_name][page_num] = []

                    if rel_url not in manifest[file_name][page_num]:
                        manifest[file_name][page_num].append(rel_url)

                    img_idx += 1
                    total_extracted += 1
                except Exception as ex:
                    continue

    return total_extracted

def run_extraction():
    print("==========================================================")
    print(" Extracting diagrams and figures from SRM source documents")
    print("==========================================================")

    manifest = {}
    total_docs = 0
    total_images = 0

    for root, dirs, files in os.walk(SOURCE_DIR):
        for f in files:
            if f.startswith('.'):
                continue
            ext = os.path.splitext(f)[1].lower()
            file_path = os.path.join(root, f)

            if ext == '.pdf':
                count = extract_pdf_images(file_path, f, manifest)
                if count > 0:
                    total_images += count
                    total_docs += 1
                    print(f"[{total_docs}] Extracted {count} images from PDF: {f}")
            elif ext == '.pptx':
                count = extract_pptx_images(file_path, f, manifest)
                if count > 0:
                    total_images += count
                    total_docs += 1
                    print(f"[{total_docs}] Extracted {count} images from PPTX: {f}")

    with open(MANIFEST_PATH, 'w', encoding='utf-8') as f:
        json.dump(manifest, f, indent=2)

    print("==========================================================")
    print(f"Extraction Complete: {total_images} images extracted from {total_docs} documents.")
    print(f"Manifest written to {MANIFEST_PATH}")
    print("==========================================================")

if __name__ == "__main__":
    run_extraction()
